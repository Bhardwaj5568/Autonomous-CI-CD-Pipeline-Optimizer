from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT = "Autonomous_CI_CD_Pipeline_Optimizer_Presentation.pptx"

TITLE_COLOR = RGBColor(15, 23, 42)
ACCENT_BLUE = RGBColor(37, 99, 235)
TEXT_COLOR = RGBColor(31, 41, 55)
BG_COLOR = RGBColor(248, 250, 252)

slides_data = [
    (
        "Autonomous CI/CD Pipeline Optimizer",
        [
            "Faster, safer, and more predictable software delivery",
            "Integration-first platform for GitHub Actions, GitLab CI, and Jenkins",
        ],
    ),
    (
        "Business Problem",
        [
            "Slow and flaky pipelines delay feature releases",
            "Risky releases cause rollbacks and outages",
            "Teams lack clear deploy/go-no-go signal",
            "CI/CD visibility is fragmented across tools",
        ],
    ),
    (
        "Solution Overview",
        [
            "Ingests CI/CD events from multiple sources",
            "Normalizes data into a common schema",
            "Calculates release risk score",
            "Recommends action: Deploy / Canary / Delay / Block",
            "Provides live operational status UI and KPIs",
        ],
    ),
    (
        "Core Features Implemented",
        [
            "Canonical event ingestion pipeline",
            "Source connectors for GitHub Actions, GitLab CI, Jenkins",
            "Queue-based asynchronous processing",
            "Rule-based risk scoring engine",
            "RBAC-style access checks and optional API key",
            "Audit logs and feedback loop",
            "KPI and status endpoints",
        ],
    ),
    (
        "Architecture (High-Level)",
        [
            "1. CI/CD source emits run data",
            "2. Connector maps payload to normalized schema",
            "3. Events enter async queue",
            "4. Scoring engine computes risk and recommendation",
            "5. API/UI exposes checks, KPIs, and audit trail",
        ],
    ),
    (
        "Security and Governance",
        [
            "Role-based header access (viewer/operator/admin)",
            "Optional API key control",
            "Audit trail for system actions",
            "Safe sharing practices and non-secret configuration model",
            "No secret values hardcoded in repo",
        ],
    ),
    (
        "Live PASS/FAIL Operations View",
        [
            "Blue indicator = PASS",
            "Red indicator = FAIL",
            "Overall status banner: OVERALL PASS / OVERALL FAIL",
            "Auto-refresh for live operational signal",
        ],
    ),
    (
        "Demo Scenarios",
        [
            "PASS demo: valid webhook payload, all_passed=true",
            "FAIL demo: malformed payload, all_passed=false",
            "Queue error check fails in FAIL scenario",
        ],
    ),
    (
        "Validation Evidence",
        [
            "Unit tests passed",
            "API health checks successful",
            "PASS and FAIL scripts verified",
            "Status UI and /status/checks verified",
            "GitHub repository updated with latest docs and code",
        ],
    ),
    (
        "Business Value",
        [
            "Faster release decisions",
            "Reduced rollback risk",
            "Higher deployment confidence",
            "Better cross-tool observability",
            "Improved engineering governance and reporting",
        ],
    ),
    (
        "Current Scope and Known Limits",
        [
            "Local runnable backend MVP",
            "Rule-based risk engine and multi-source ingestion",
            "Queue fail counter resets on restart",
            "Status UI is operational snapshot, not long-term BI",
        ],
    ),
    (
        "Roadmap",
        [
            "Jenkins production onboarding flow",
            "Persistent queue and production datastore hardening",
            "CI automation for build/test/release",
            "Advanced risk calibration from real workload history",
        ],
    ),
]

prs = Presentation()

for title, bullets in slides_data:
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # Background tint
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    title_shape = slide.shapes.title
    title_shape.text = title
    title_tf = title_shape.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.font.size = Pt(34)
    title_p.font.bold = True
    title_p.font.color.rgb = TITLE_COLOR

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_COLOR

    # Accent bar on left
    bar = slide.shapes.add_shape(
        autoshape_type_id=1,  # rectangle
        left=Inches(0.2),
        top=Inches(0.0),
        width=Inches(0.08),
        height=Inches(7.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

prs.save(OUTPUT)
print(OUTPUT)
